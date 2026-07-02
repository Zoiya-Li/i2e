"""Offline tests for the diagram2ppt v2 iterative pipeline (no network).

A MockVLM replays scripted responses; originals are synthesized with the v2
renderer itself, so a correct IR scores ~zero residual by construction and
every failure injected below is unambiguous.

Run: python tests/test_diagram_v2.py   (or pytest tests/test_diagram_v2.py)
"""
from __future__ import annotations

import json
import os
import sys
import tempfile

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import numpy as np
from PIL import Image

from work.diagram2ppt.v2 import diff as diff_mod
from work.diagram2ppt.v2 import ir as ir_mod
from work.diagram2ppt.v2 import parsing
from work.diagram2ppt.v2.build_pptx import build_pptx
from work.diagram2ppt.v2.loop import run_loop, _padded_crop
from work.diagram2ppt.v2.render import render

W, H = 800, 400


class MockVLM:
    """Replays scripted responses; entries may be str or (prompt, image)->str."""
    CROP_MAX_EDGE = 1024

    def __init__(self, responses):
        self.responses = list(responses)
        self.calls = 0

    def chat(self, prompt, image, max_edge=None):
        self.calls += 1
        if not self.responses:
            raise AssertionError("MockVLM ran out of scripted responses")
        r = self.responses.pop(0)
        return r(prompt, image) if callable(r) else r


def _vlm_json(elements: list[dict]) -> str:
    return json.dumps({"elements": elements})


# Two filled rects + a connecting arrow + a free text block, as the VLM
# would describe them (fractions of an 800x400 canvas).
RECT_A = {"id": "a", "type": "rect", "x": 0.0625, "y": 0.25, "width": 0.25,
          "height": 0.25, "text": "Input", "fill": "#4A90D9",
          "border_color": "#2C5F8A", "text_color": "#FFFFFF"}
RECT_B = {"id": "b", "type": "rect", "x": 0.625, "y": 0.25, "width": 0.25,
          "height": 0.25, "text": "Output", "fill": "#4CAF50",
          "border_color": "#2E7D32", "text_color": "#FFFFFF"}
ARROW = {"id": "c", "type": "arrow", "from_id": "a", "to_id": "b",
         "color": "#333333"}
TEXT = {"id": "t", "type": "text", "x": 0.3, "y": 0.75, "width": 0.4,
        "height": 0.1, "text": "Pipeline", "text_color": "#000000"}
SCENE = [RECT_A, RECT_B, ARROW, TEXT]


def _original(elements: list[dict]) -> tuple[Image.Image, dict]:
    """Synthesize the 'original' image by rendering a ground-truth IR."""
    truth = ir_mod.new_ir("synthetic", W, H)
    truth["elements"] = ir_mod.from_vlm_elements(elements, W, H)
    return render(truth), truth


# ---------------------------------------------------------------- parsing --

def test_parsing_variants():
    els = [{"id": "x", "type": "rect", "x": 0.1, "y": 0.1,
            "width": 0.2, "height": 0.2}]
    obj = json.dumps({"elements": els})
    arr = json.dumps(els)
    assert parsing.parse_elements(obj) == els
    assert parsing.parse_elements(arr) == els
    assert parsing.parse_elements(f"```json\n{obj}\n```") == els
    assert parsing.parse_elements(f"Here is the JSON you asked for:\n{obj}\nDone!") == els
    assert parsing.parse_elements(f"JSON\n{arr} -- that's all") == els
    try:
        parsing.parse_elements("I see a nice diagram with boxes.")
        raise AssertionError("expected ValueError on non-JSON")
    except ValueError:
        pass

    # max_tokens truncation: salvage the complete leading objects
    el2 = dict(els[0], id="y")
    truncated = json.dumps({"elements": [els[0], el2]})[:-30]  # cut mid-object
    salvaged = parsing.parse_elements(truncated)
    assert salvaged == [els[0]], salvaged
    print("ok parsing variants (incl. truncation salvage)")


# ----------------------------------------------------- render-diff signal --

def test_residual_separates_right_from_wrong():
    original, truth = _original(SCENE)
    rendered = render(truth)
    a = truth["elements"][0]
    good = diff_mod.element_residual(original, rendered, a["bbox"])
    assert good < 0.05, f"identical render should be ~0, got {good}"

    broken = json.loads(json.dumps(truth))  # deep copy
    broken["elements"][0]["fill"] = ""      # lose the solid fill
    bad = diff_mod.element_residual(original, render(broken), a["bbox"])
    assert bad > 0.45, f"missing fill should score high, got {bad}"

    # sliver bboxes must not crash SSIM (win_size>=7 constraint)
    for sliver in ([10, 10, 14, 200], [10, 10, 200, 13], [5, 5, 7, 7]):
        r = diff_mod.element_residual(original, rendered, sliver)
        assert 0.0 <= r <= 1.0, (sliver, r)
    print(f"ok residual signal (good={good}, bad={bad})")


def test_coverage_finds_missed_shape_and_covers_connectors():
    original, truth = _original(SCENE)
    cov_full = diff_mod.coverage(original, truth)
    assert cov_full["explained_frac"] > 0.97, cov_full  # arrow ink covered too

    partial = ir_mod.new_ir("synthetic", W, H)
    partial["elements"] = ir_mod.from_vlm_elements([RECT_A], W, H)
    cov = diff_mod.coverage(original, partial)
    assert cov["explained_frac"] < 0.8, cov
    assert cov["missing"], "should report missed regions"
    mx0, my0, mx1, my1 = cov["missing"][0]["bbox"]
    bx0, by0, bx1, by1 = ir_mod.from_vlm_elements([RECT_B], W, H)[0]["bbox"]
    assert mx0 < bx1 and mx1 > bx0 and my0 < by1 and my1 > by0, \
        "largest missing region should overlap the missed rect"
    print(f"ok coverage (full={cov_full['explained_frac']}, partial={cov['explained_frac']})")


def test_text_residual_font_tolerant_but_catches_missing():
    original, truth = _original(SCENE)
    t = next(e for e in truth["elements"] if e["type"] == "text")

    # same text, very different size (proxy for a font difference) → keep
    styled = json.loads(json.dumps(truth))
    st = next(e for e in styled["elements"] if e["type"] == "text")
    st["font_size"] = 14.0
    r_font = diff_mod.text_residual(original, render(styled), t["bbox"])
    assert r_font < 0.62, f"font/style difference must not demote, got {r_font}"

    # text missing entirely → demote
    gone = json.loads(json.dumps(truth))
    gone["elements"] = [e for e in gone["elements"] if e["type"] != "text"]
    r_gone = diff_mod.text_residual(original, render(gone), t["bbox"])
    assert r_gone > 0.62, f"missing text must score high, got {r_gone}"
    print(f"ok text residual (font={r_font}, missing={r_gone})")


def test_shell_scoring_ignores_children():
    container = {"id": "panel", "type": "rect", "x": 0.05, "y": 0.1,
                 "width": 0.5, "height": 0.7, "fill": "#EEF2F8",
                 "border_color": "#2C5F8A"}
    child = {"id": "inner", "type": "rect", "x": 0.1, "y": 0.2, "width": 0.4,
             "height": 0.5, "fill": "#CC2222"}
    original, truth = _original([container, child])
    panel = truth["elements"][0]

    # child garbled in the render — lost its fill entirely; shell unchanged
    broken = json.loads(json.dumps(truth))
    broken["elements"][1]["fill"] = ""
    rendered = render(broken)

    naive = diff_mod.element_residual(original, rendered, panel["bbox"])
    kids = diff_mod.children_of(panel, broken["elements"])
    shell = diff_mod.element_residual(original, rendered, panel["bbox"],
                                      exclude=kids)
    assert kids, "child bbox should be detected as contained"
    assert shell < 0.2 < naive, (shell, naive)
    print(f"ok shell scoring (naive={naive}, shell={shell})")


def test_faithful_crop_punches_native_children():
    from work.diagram2ppt.v2.render import faithful_crop
    container = {"id": "panel", "type": "rect", "x": 0.1, "y": 0.1,
                 "width": 0.5, "height": 0.6, "fill": "#3355AA"}
    child = {"id": "inner", "type": "text", "x": 0.2, "y": 0.25, "width": 0.3,
             "height": 0.2, "text": "EDIT ME", "text_color": "#FFFFFF"}
    original, truth = _original([container, child])
    panel, inner = truth["elements"]
    ir_mod.demote(panel)  # panel ships as a screenshot; child stays native

    crop = faithful_crop(original, panel, truth["elements"])
    px0, py0 = panel["bbox"][0], panel["bbox"][1]
    cx0, cy0, cx1, cy1 = (int(v) for v in inner["bbox"])
    region = np.asarray(crop)[cy0 - int(py0):cy1 - int(py0),
                              cx0 - int(px0):cx1 - int(px0)]
    # the white glyph pixels must be patched to the panel's flat fill
    assert region.std(axis=(0, 1)).max() < 1.0, "child region should be flat"
    assert abs(int(region[..., 2].mean()) - 0xAA) < 8, "patched with panel fill"
    print("ok faithful crop punch")


# ------------------------------------------------------------------- loop --

def test_loop_converges_on_correct_extraction():
    original, _ = _original(SCENE)
    with tempfile.TemporaryDirectory() as td:
        src = os.path.join(td, "orig.png")
        original.save(src)
        vlm = MockVLM([_vlm_json(SCENE)])
        result = run_loop(src, vlm, td, max_rounds=3, log=lambda *_: None)
    m = result["history"][-1]
    assert vlm.calls == 1, f"correct extraction needs only the global call, used {vlm.calls}"
    assert m["native_fraction_count"] == 1.0, m
    assert m["demoted_count"] == 0, m
    assert m["coverage"] > 0.97, m
    print(f"ok loop converges ({json.dumps(m)})")


def test_loop_refine_fixes_bad_element():
    original, truth = _original(SCENE)
    bad_a = dict(RECT_A, fill="", border_color="#2C5F8A")  # lost its fill

    # refine sees a crop around the bad bbox; answer with the truth mapped
    # into crop fractions (same mapping the real VLM would express).
    true_bbox = truth["elements"][0]["bbox"]
    pred_bbox = ir_mod.from_vlm_elements([bad_a], W, H)[0]["bbox"]
    _, crop_box = _padded_crop(original, pred_bbox)
    cw, ch = crop_box[2] - crop_box[0], crop_box[3] - crop_box[1]
    fixed = dict(RECT_A,
                 x=(true_bbox[0] - crop_box[0]) / cw,
                 y=(true_bbox[1] - crop_box[1]) / ch,
                 width=(true_bbox[2] - true_bbox[0]) / cw,
                 height=(true_bbox[3] - true_bbox[1]) / ch)

    with tempfile.TemporaryDirectory() as td:
        src = os.path.join(td, "orig.png")
        original.save(src)
        vlm = MockVLM([_vlm_json([bad_a, RECT_B, ARROW, TEXT]),
                       _vlm_json([fixed])])
        result = run_loop(src, vlm, td, max_rounds=3,
                          coverage_target=0.0,  # isolate the refine path
                          log=lambda *_: None)
    m = result["history"][-1]
    a = next(e for e in result["elements"] if e["id"] == "a")
    assert a["status"] == "native" and a["fill"] == "#4A90D9", a
    assert a["residual"] < 0.1, a["residual"]
    assert m["demoted_count"] == 0, m
    print(f"ok refine path (residual={a['residual']})")


def test_loop_demotes_stubborn_element_and_stays_faithful():
    original, _ = _original(SCENE)
    bad_a = dict(RECT_A, fill="")
    with tempfile.TemporaryDirectory() as td:
        src = os.path.join(td, "orig.png")
        original.save(src)
        # refine keeps returning the same broken element — never converges
        vlm = MockVLM([_vlm_json([bad_a, RECT_B, ARROW, TEXT]),
                       _vlm_json([bad_a]), _vlm_json([bad_a])])
        result = run_loop(src, vlm, td, max_rounds=4, max_tries=2,
                          coverage_target=0.0, log=lambda *_: None)
        a = next(e for e in result["elements"] if e["id"] == "a")
        assert a["type"] == "raster_crop" and a["status"] == "demoted", a
        assert a["ext"]["original_type"] == "rect", a

        # the fidelity guarantee: inside the demoted bbox, the re-render IS
        # the original
        rerender = render(result, original)
        x0, y0, x1, y1 = (int(v) for v in a["bbox"])
        diff = np.abs(
            np.asarray(rerender.crop((x0, y0, x1, y1)), dtype=np.int16)
            - np.asarray(original.crop((x0, y0, x1, y1)), dtype=np.int16))
        assert diff.max() == 0, f"demoted crop must be pixel-identical, max diff {diff.max()}"
    print("ok demote path (pixel-faithful)")


def test_loop_unidentifiable_region_falls_back_to_crop():
    original, _ = _original(SCENE)
    with tempfile.TemporaryDirectory() as td:
        src = os.path.join(td, "orig.png")
        original.save(src)
        # global misses rect B; identify answers garbage → crop fallback
        vlm = MockVLM([_vlm_json([RECT_A, TEXT]),
                       "I cannot make out any structure here."] + ["{}"] * 6)
        result = run_loop(src, vlm, td, max_rounds=3, log=lambda *_: None)
    crops = [e for e in result["elements"] if e["type"] == "raster_crop"]
    assert crops, "missed ink must come back as faithful crops"
    bx0, by0, bx1, by1 = ir_mod.from_vlm_elements([RECT_B], W, H)[0]["bbox"]
    assert any(c["bbox"][0] < bx1 and c["bbox"][2] > bx0 and
               c["bbox"][1] < by1 and c["bbox"][3] > by0 for c in crops), \
        "a crop should land on the missed rect"
    assert result["history"][-1]["coverage"] > 0.9, result["history"][-1]
    print("ok identify fallback")


def test_ocr_snap_geometry_and_additions():
    from work.diagram2ppt.v2.ocr_snap import snap_text
    ir = ir_mod.new_ir("synthetic", W, H)
    ir["elements"] = ir_mod.from_vlm_elements([
        dict(TEXT, text="Pipeline Overview\nSecond Line"),   # multi-line, loose bbox
        {"id": "ras", "type": "raster", "x": 0.7, "y": 0.6,
         "width": 0.25, "height": 0.3},                      # chart region
    ], W, H)

    lines = [
        {"bbox": [240, 300, 410, 322], "text": "Pipeline 0verview", "conf": 0.93},  # OCR'd l→0
        {"bbox": [250, 330, 400, 350], "text": "Second Line", "conf": 0.97},
        {"bbox": [100, 40, 220, 70], "text": "Missed Header", "conf": 0.95},        # unclaimed
        {"bbox": [580, 300, 640, 315], "text": "axis label", "conf": 0.92},         # inside raster
        {"bbox": [60, 200, 90, 215], "text": "x?", "conf": 0.4},                    # low conf
    ]
    r = snap_text(ir, lines, log=lambda *_: None)
    assert r == {"snapped": 1, "added": 1}, r

    t = ir["elements"][0]
    assert t["bbox"] == [240.0, 300.0, 410.0, 350.0], t["bbox"]  # union of 2 lines
    assert t["text"] == "Pipeline Overview\nSecond Line"          # VLM content kept
    assert t["ext"]["ocr"] == "snap"
    assert 15 < t["font_size"] < 20, t["font_size"]               # median line height *0.82

    added = [e for e in ir["elements"] if e.get("ext", {}).get("ocr") == "added"]
    assert len(added) == 1 and added[0]["text"] == "Missed Header", added
    print("ok ocr snap")


def test_fit_font_respects_width():
    from work.diagram2ppt.v2.build_pptx import _fit_font_px
    # tall narrow box: height alone would say ~86px; width must cap it
    tall = {"bbox": [0, 0, 60, 120], "text": "RETAIN", "font_size": None}
    size = _fit_font_px(tall)
    # width_factor is 0.53; the cap is slightly higher than the old 0.55
    assert size <= 60 / 6 / 0.53 + 0.01, size
    assert size < 30, f"narrow box must not get a giant font, got {size}"
    # wide short box: height is the binding constraint
    # single-line text uses fit_height_factor 0.90, multi-line uses 0.72
    wide = {"bbox": [0, 0, 800, 30], "text": "short", "font_size": None}
    assert 15 < _fit_font_px(wide) <= 30 * 0.90 + 0.01
    # explicit font_size only caps, never inflates
    capped = {"bbox": [0, 0, 800, 30], "text": "short", "font_size": 12}
    assert _fit_font_px(capped) == 12
    print("ok fit font")


def test_dedupe_text_drops_shape_label_duplicates():
    from work.diagram2ppt.v2.ocr_snap import dedupe_text
    ir = ir_mod.new_ir("synthetic", W, H)
    ir["elements"] = ir_mod.from_vlm_elements([
        RECT_A,                                                   # label "Input"
        {"id": "dup", "type": "text", "x": 0.08, "y": 0.3,
         "width": 0.2, "height": 0.1, "text": "Input"},           # overlaps A
        {"id": "far", "type": "text", "x": 0.7, "y": 0.8,
         "width": 0.2, "height": 0.1, "text": "Input"},           # elsewhere: keep
    ], W, H)
    n = dedupe_text(ir, log=lambda *_: None)
    ids = [e["id"] for e in ir["elements"]]
    assert n == 1 and "dup" not in ids and "far" in ids, ids
    print("ok dedupe text")


def test_prune_connectors_keeps_real_drops_hallucinated():
    original, truth = _original(SCENE)
    truth["elements"].append({               # no ink between rect A and TEXT
        "id": "ghost", "type": "arrow", "from_id": "a", "to_id": "t",
        "color": "#333333", "status": "native", "tries": 0,
        "residual": None, "z": 9, "ext": {},
    })
    dropped = diff_mod.prune_connectors(original, truth, log=lambda *_: None)
    ids = [e["id"] for e in truth["elements"]]
    assert dropped == 1 and "ghost" not in ids and "c" in ids, ids
    print("ok prune connectors")


# ------------------------------------------------------------------- pptx --

def test_build_pptx_native_and_crops():
    original, truth = _original(SCENE)
    with tempfile.TemporaryDirectory() as td:
        src = os.path.join(td, "orig.png")
        original.save(src)
        truth["image"]["path"] = src
        ir_mod.demote(truth["elements"][1])  # ship rect B as a picture
        out = os.path.join(td, "deck.pptx")
        counts = build_pptx(truth, out)

        assert counts == {"shapes": 2, "pictures": 1, "connectors": 1}, counts
        from pptx import Presentation
        prs = Presentation(out)
        slide = prs.slides[0]
        assert len(slide.shapes) == 4, len(slide.shapes)
        texts = {sh.text_frame.text for sh in slide.shapes if sh.has_text_frame}
        assert "Input" in texts and "Pipeline" in texts, texts
    print(f"ok pptx build ({json.dumps(counts)})")


def test_experts_formula_and_chart_pptx():
    from work.diagram2ppt.v2 import experts

    ir = ir_mod.new_ir("synthetic", W, H)
    ir["elements"] = ir_mod.from_vlm_elements([RECT_A], W, H)
    ir["elements"].append({
        "id": "f1", "type": "formula", "status": "native", "tries": 0,
        "residual": None, "z": 5, "bbox": [400.0, 50.0, 640.0, 110.0],
        "text": "A = ...", "latex": "A \\approx 1", "fill": "",
        "border_color": "", "text_color": "", "bold": False, "font_size": None,
        "omml": '<m:oMath><m:r><m:t>A</m:t></m:r><m:r><m:t>≈1</m:t></m:r></m:oMath>',
        "ext": {"expert": "formula"},
    })
    ir["elements"].append({
        "id": "f2", "type": "formula", "status": "native", "tries": 0,
        "residual": None, "z": 6, "bbox": [400.0, 150.0, 640.0, 200.0],
        "text": "", "latex": "\\beta = 2", "fill": "", "border_color": "",
        "text_color": "", "bold": False, "font_size": None,
        "ext": {"expert": "formula"},   # no omml → latex text fallback
    })
    ir["elements"].append({
        "id": "c1", "type": "chart", "status": "native", "tries": 0,
        "residual": None, "z": 7, "bbox": [80.0, 220.0, 380.0, 380.0],
        "text": "", "fill": "", "border_color": "", "text_color": "",
        "bold": False, "font_size": None,
        "chart": {"kind": "bar", "categories": ["orthogonal", "aligned"],
                  "series": [{"name": "coverage", "color": "#4472C4",
                              "values": [0.78, 0.12]}]},
        "ext": {"expert": "chart", "approx": True},
    })

    # candidates: garbled math text is detected, plain text is not
    ir["elements"].append(dict(ir_mod.from_vlm_elements(
        [dict(TEXT, id="mathy", text="θ ≈ 0°")], W, H)[0], id="mathy"))
    cands = experts.formula_candidates(ir)
    assert any(e["id"] == "mathy" for e in cands), "math text should be flagged"
    assert all(e["id"] != "t" for e in cands), "plain text must not be flagged"

    with tempfile.TemporaryDirectory() as td:
        out = os.path.join(td, "deck.pptx")
        counts = build_pptx(ir, out)
        assert counts["charts"] == 1 and counts["shapes"] >= 3, counts
        import zipfile
        z = zipfile.ZipFile(out)
        slide_xml = z.read("ppt/slides/slide1.xml").decode("utf-8")
        assert "oMath" in slide_xml, "OMML equation must be embedded"
        # LaTeX fallback is now converted to editable OMML unicode, not raw source
        assert "β = 2" in slide_xml, "no-OMML formula should be converted to OMML unicode"
        assert any(n.startswith("ppt/charts/") for n in z.namelist()), \
            "native chart part must exist"
    print("ok experts formula+chart pptx")


def test_sanitize_drops_slivers():
    from work.diagram2ppt.v2.experts import sanitize
    ir = ir_mod.new_ir("synthetic", W, H)
    ir["elements"] = [
        {"id": "ok", "type": "raster_crop", "status": "demoted", "tries": 0,
         "residual": 0.0, "z": 1, "bbox": [10, 10, 60, 60], "ext": {}},
        {"id": "strip", "type": "raster_crop", "status": "demoted", "tries": 0,
         "residual": 0.0, "z": 2, "bbox": [0, 100, 500, 103], "ext": {}},
        {"id": "zero", "type": "raster_crop", "status": "demoted", "tries": 0,
         "residual": 0.0, "z": 3, "bbox": [700, 10, 700, 47], "ext": {}},
    ]
    n = sanitize(ir, log=lambda *_: None)
    ids = [e["id"] for e in ir["elements"]]
    assert n == 2 and ids == ["ok"], ids
    print("ok sanitize")


def test_vectorize_dots_and_silhouette():
    from work.diagram2ppt.v2 import vectorize
    from PIL import ImageDraw

    crop = Image.new("RGB", (200, 150), "#F4F6FA")
    d = ImageDraw.Draw(crop)
    truth = [(40, 40, "#cc2222"), (100, 70, "#2255cc"), (160, 110, "#22aa44")]
    for cx, cy, col in truth:
        d.ellipse([cx - 4, cy - 4, cx + 4, cy + 4], fill=col)
    dots = vectorize.extract_dots(crop)
    assert len(dots) == 3, dots
    for cx, cy, col in truth:
        hit = min(dots, key=lambda dd: abs(dd["cx"] - cx) + abs(dd["cy"] - cy))
        assert abs(hit["cx"] - cx) < 3 and abs(hit["cy"] - cy) < 3, (hit, cx, cy)
        want = tuple(int(col[i:i + 2], 16) for i in (1, 3, 5))
        got = tuple(int(hit["color"][i:i + 2], 16) for i in (1, 3, 5))
        assert max(abs(a - b) for a, b in zip(want, got)) < 40, (want, got)

    blob = Image.new("RGB", (200, 150), "white")
    ImageDraw.Draw(blob).ellipse([30, 30, 170, 120], fill="#88aacc")
    sil = vectorize.extract_silhouette(blob)
    assert sil and len(sil["points"]) >= 3
    xs = [p[0] for p in sil["points"]]; ys = [p[1] for p in sil["points"]]
    assert 20 < min(xs) < 45 and 155 < max(xs) < 180, (min(xs), max(xs))
    print(f"ok vectorize (3 dots recovered, silhouette {len(sil['points'])} pts)")


def test_ban_screenshots_build_has_zero_pictures():
    from work.diagram2ppt.v2.vectorize import vectorize_pass

    original, truth = _original(SCENE)
    with tempfile.TemporaryDirectory() as td:
        src = os.path.join(td, "orig.png")
        original.save(src)
        truth["image"]["path"] = src
        # one failed shape + one organic crop + one icon-sized crop
        ir_mod.demote(truth["elements"][1])                       # rect B
        truth["elements"].append({
            "id": "art", "type": "raster_crop", "status": "demoted",
            "tries": 0, "residual": 0.0, "z": 8,
            "bbox": [80.0, 250.0, 380.0, 390.0], "ext": {"original_type": "raster"},
        })
        truth["elements"].append({
            "id": "ic", "type": "raster_crop", "status": "demoted",
            "tries": 0, "residual": 0.0, "z": 9,
            "bbox": [600.0, 300.0, 650.0, 350.0], "ext": {"original_type": "raster"},
        })

        class IconVLM:
            CROP_MAX_EDGE = 1024
            calls = 0
            def chat(self, prompt, image, max_edge=None):
                self.calls += 1
                return '{"kind": "database", "color": "#4472C4", "glyph": "🗄"}'

        stats = vectorize_pass(truth, original, IconVLM(), log=lambda *_: None)
        assert stats["forced_shapes"] == 1 and stats["icons"] == 1, stats
        assert all(e["type"] != "raster_crop" for e in truth["elements"])

        out = os.path.join(td, "deck.pptx")
        counts = build_pptx(truth, out)
        assert counts["pictures"] == 0, counts
        assert counts.get("icons") == 1, counts
        from pptx import Presentation
        assert len(Presentation(out).slides[0].shapes) >= 4
    print(f"ok ban-screenshots build ({json.dumps(counts)})")


ALL = [v for k, v in sorted(globals().items()) if k.startswith("test_")]

if __name__ == "__main__":
    for fn in ALL:
        fn()
    print(f"\n{len(ALL)}/{len(ALL)} diagram2ppt-v2 tests green")
