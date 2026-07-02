"""Deterministic structural stats for a ``.pptx`` — no model, no network.

Used to (a) lock the frozen v2 delivery as a regression baseline and (b) score
v3 output against it: native-object ratio, picture (raster fallback) count, and
OMML formula count. Everything here is derived straight from the file bytes so
it is fully reproducible offline.
"""
from __future__ import annotations

import collections
import hashlib
import re
import zipfile
from pathlib import Path
from typing import Any

# python-pptx MSO_SHAPE_TYPE integer values.
_GROUP = 6
_PICTURE = 13

# Shape types that count as native / editable objects (everything except raster
# PICTUREs). PLACEHOLDER/GROUP are structural and excluded from the ratio.
_NATIVE_TYPES = {"AUTO_SHAPE", "TEXT_BOX", "LINE", "FREEFORM", "CHART", "TABLE"}


def _walk(shapes):
    for sh in shapes:
        yield sh
        if int(getattr(sh, "shape_type", -1) or -1) == _GROUP:
            yield from _walk(sh.shapes)


def pptx_structure(path: str | Path) -> dict[str, Any]:
    """Return a deterministic structural fingerprint of a presentation."""
    from pptx import Presentation

    path = Path(path)
    prs = Presentation(str(path))
    hist: collections.Counter = collections.Counter()
    top = 0
    total = 0
    pictures = 0
    for slide in prs.slides:
        top += len(slide.shapes)
        for sh in _walk(slide.shapes):
            total += 1
            name = str(getattr(sh, "shape_type", "UNKNOWN")).split(" ")[0]
            hist[name] += 1
            if int(getattr(sh, "shape_type", -1) or -1) == _PICTURE:
                pictures += 1

    omml = 0
    with zipfile.ZipFile(str(path)) as z:
        for name in z.namelist():
            if name.endswith(".xml") and "slide" in name:
                xml = z.read(name).decode("utf-8", "ignore")
                omml += len(re.findall(r"<a14:m>|<m:oMath\b", xml))

    native = sum(c for t, c in hist.items() if t in _NATIVE_TYPES)
    native_object_ratio = round(native / total, 4) if total else 0.0
    return {
        "slides": len(prs.slides),
        "top_level_shapes": top,
        "total_shapes_recursive": total,
        "shape_histogram": dict(sorted(hist.items())),
        "pictures": pictures,
        "omml_math_runs": omml,
        "native_object_ratio": native_object_ratio,
    }


def sha256_file(path: str | Path) -> str:
    return hashlib.sha256(Path(path).read_bytes()).hexdigest()
