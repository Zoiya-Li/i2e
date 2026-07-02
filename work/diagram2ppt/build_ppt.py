"""Step 2: Build native PPTX from structured diagram JSON.

Each element type maps to a native PowerPoint shape:
  - rect/rounded_rect → Rectangle / RoundedRectangle
  - oval/circle       → Oval
  - diamond           → Diamond
  - text              → Rectangle with no fill/border (text only)
  - arrow             → Connector with arrowhead
  - line              → Connector (no arrowhead)

All positions are in EMU (English Metric Units), derived from
fraction coordinates in the analysis JSON.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn


# Standard 16:9 slide
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

SHAPE_TYPE_MAP = {
    "rect":           MSO_SHAPE.RECTANGLE,
    "rounded_rect":   MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval":           MSO_SHAPE.OVAL,
    "circle":         MSO_SHAPE.OVAL,
    "diamond":        MSO_SHAPE.DIAMOND,
    "parallelogram":  MSO_SHAPE.PARALLELOGRAM,
    "hexagon":        MSO_SHAPE.HEXAGON,
}


def build_ppt(data: dict, output_path: str,
              image_w: int = None, image_h: int = None):
    """Build a native PPTX from structured diagram data.

    Args:
        data: {"elements": [...]} from analyze.py
        output_path: path to save .pptx
        image_w, image_h: original image dimensions (for aspect ratio)
    """
    prs = Presentation()

    # Match aspect ratio to source image
    if image_w and image_h:
        ratio = image_w / image_h
        target_w = SLIDE_W
        target_h = int(SLIDE_W / ratio)
        if target_h > SLIDE_H:
            target_h = SLIDE_H
            target_w = int(SLIDE_H * ratio)
        prs.slide_width = target_w
        prs.slide_height = target_h
    else:
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

    sw = prs.slide_width
    sh = prs.slide_height

    # Blank layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    shape_map = {}  # id → shape, for connectors

    # First pass: shapes (skip connectors)
    for el in data.get("elements", []):
        if el.get("type") in ("arrow", "line"):
            continue
        shape = _add_shape(slide, el, sw, sh)
        if shape and el.get("id"):
            shape_map[el["id"]] = shape

    # Second pass: connectors (need shape_map to resolve IDs)
    for el in data.get("elements", []):
        if el.get("type") not in ("arrow", "line"):
            continue
        _add_connector(slide, el, shape_map)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    w_in = sw / 914400
    h_in = sh / 914400
    print(f"✓ PPT saved → {output_path} ({w_in:.1f}\" × {h_in:.1f}\")")


# ---------------------------------------------------------------------------
# Shape helpers
# ---------------------------------------------------------------------------

def _hex_to_rgb(hex_str: str):
    """Parse '#RRGGBB' to RGBColor, or return None."""
    if not hex_str or hex_str in ("none", "transparent", "null", "None"):
        return None
    hex_str = hex_str.lstrip('#')
    if len(hex_str) != 6:
        return None
    try:
        return RGBColor(int(hex_str[0:2], 16),
                        int(hex_str[2:4], 16),
                        int(hex_str[4:6], 16))
    except (ValueError, TypeError):
        return None


def _add_shape(slide, el: dict, sw: int, sh: int):
    """Create a native PPT shape for one element."""
    el_type = el.get("type", "rect")

    # Map element type → PPT shape type
    if el_type == "text":
        ppt_type = MSO_SHAPE.RECTANGLE
    elif el_type in ("icon", "badge"):
        ppt_type = MSO_SHAPE.OVAL
    else:
        ppt_type = SHAPE_TYPE_MAP.get(el_type, MSO_SHAPE.RECTANGLE)

    # Position: fraction → EMU
    x = int(el.get("x", 0) * sw)
    y = int(el.get("y", 0) * sh)
    w = int(el.get("width", 0.1) * sw)
    h = int(el.get("height", 0.05) * sh)

    shape = slide.shapes.add_shape(ppt_type, x, y, w, h)

    # --- Fill ---
    fill_color = _hex_to_rgb(el.get("fill"))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    # --- Border ---
    border_color = _hex_to_rgb(el.get("border_color"))
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(el.get("border_width", 1))
    else:
        shape.line.fill.background()

    # --- Text ---
    text = el.get("text", "")
    if text:
        tf = shape.text_frame
        tf.word_wrap = True

        # Clear default paragraph
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER

        run = p.add_run()
        run.text = text

        # Text color
        tc = _hex_to_rgb(el.get("text_color"))
        if tc:
            run.font.color.rgb = tc

        # Bold
        if el.get("bold"):
            run.font.bold = True

        # Font size
        fs = el.get("font_size")
        if fs:
            run.font.size = Pt(fs)

        # Vertical centering (approximate via margins)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)

    return shape


# ---------------------------------------------------------------------------
# Connector helpers
# ---------------------------------------------------------------------------

def _shape_center(shape) -> tuple:
    return (shape.left + shape.width // 2,
            shape.top + shape.height // 2)


def _edge_point(shape, target_center: tuple) -> tuple:
    """Point on shape's bounding box closest to target_center."""
    cx, cy = _shape_center(shape)
    tx, ty = target_center
    dx, dy = tx - cx, ty - cy

    if dx == 0 and dy == 0:
        return (cx, cy + shape.height // 2)

    hw = shape.width / 2
    hh = shape.height / 2

    # Scale vector to touch bounding box
    sx = hw / abs(dx) if dx != 0 else float('inf')
    sy = hh / abs(dy) if dy != 0 else float('inf')
    s = min(sx, sy)

    return (int(cx + dx * s), int(cy + dy * s))


def _add_connector(slide, el: dict, shape_map: dict):
    """Draw an arrow/line connector between two shapes."""
    from_id = el.get("from_id")
    to_id = el.get("to_id")

    if not from_id or not to_id:
        return
    if from_id not in shape_map or to_id not in shape_map:
        print(f"  ⚠ Connector {el.get('id','')} skipped: "
              f"shape {from_id} or {to_id} not found")
        return

    from_shape = shape_map[from_id]
    to_shape = shape_map[to_id]

    # Edge-to-edge points
    to_center = _shape_center(to_shape)
    from_center = _shape_center(from_shape)
    start = _edge_point(from_shape, to_center)
    end = _edge_point(to_shape, from_center)

    try:
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            start[0], start[1], end[0], end[1]
        )
    except Exception as e:
        print(f"  ⚠ Connector failed: {e}")
        return

    # Color
    color = _hex_to_rgb(el.get("color"))
    if color:
        connector.line.color.rgb = color
    connector.line.width = Pt(1.5)

    # Arrowhead
    if el.get("type") == "arrow":
        _add_arrowhead(connector)

    # Label on connector
    label = el.get("label", "")
    if label:
        _add_connector_label(connector, label, color)


def _add_arrowhead(connector):
    """Add triangle arrowhead to connector end."""
    try:
        ln = connector.line._ln
        tail = ln.makeelement(qn('a:tailEnd'), {
            'type': 'triangle',
            'w': 'med',
            'len': 'med'
        })
        ln.append(tail)
    except Exception:
        pass


def _add_connector_label(connector, label: str, color):
    """Add text label at midpoint of connector."""
    try:
        from pptx.util import Emu
        cx = (connector.left + connector.width // 2)
        cy = (connector.top + connector.height // 2)
        # Offset label slightly above the line
        cy -= Inches(0.15)
        tf = connector.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = label
        run.font.size = Pt(9)
        if color:
            run.font.color.rgb = color
    except Exception:
        pass
