"""Final IR → native PPTX (pixel-bbox port of v1 build_ppt + picture crops).

Native elements become PowerPoint shapes / text boxes / connectors;
raster_crop elements become embedded pictures cut from the ORIGINAL image
(saved alongside the .pptx), so the deck recomposites faithfully even where
extraction didn't converge.
"""
from __future__ import annotations

import html
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from work.diagram2ppt.build_ppt import SHAPE_TYPE_MAP, _add_arrowhead, _hex_to_rgb
from work.diagram2ppt.v2.render import faithful_crop

SLIDE_W_IN = 13.333


def build_pptx(ir: dict, output_path: str) -> dict:
    """Build the deck. Returns {"shapes": n, "pictures": n, "connectors": n}."""
    w, h = ir["image"]["width"], ir["image"]["height"]
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_W_IN * h / w)
    scale = prs.slide_width / w                      # EMU per source px
    pt_per_px = SLIDE_W_IN * 72 / w                  # font pt per source px

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    assets = out.parent / (out.stem + "_assets")

    original = None
    counts = {"shapes": 0, "pictures": 0, "connectors": 0}
    shape_map: dict[str, object] = {}

    ordered = sorted(ir["elements"], key=lambda e: e.get("z", 0))

    for el in ordered:                               # pass 1: boxed elements
        t = el["type"]
        if t in ("arrow", "line"):
            continue
        x0, y0, x1, y1 = (int(v * scale) for v in el["bbox"])
        bw, bh = max(1, x1 - x0), max(1, y1 - y0)

        if t == "formula":
            box = slide.shapes.add_textbox(x0, y0, bw, bh)
            _set_formula(box, el, pt_per_px)
            shape_map[el["id"]] = box
            counts["shapes"] += 1
            continue

        if t == "chart":
            frame = _add_chart(slide, el, x0, y0, bw, bh)
            if frame is not None:
                shape_map[el["id"]] = frame
                counts["charts"] = counts.get("charts", 0) + 1
            continue

        if t == "icon":
            shape = _add_icon(slide, el, x0, y0, bw, bh, pt_per_px)
            shape_map[el["id"]] = shape
            counts["icons"] = counts.get("icons", 0) + 1
            continue

        if t in ("dotcloud", "surface"):
            n = _add_dotcloud(slide, el, scale)
            counts["dots"] = counts.get("dots", 0) + n
            if t == "surface":
                counts["surfaces"] = counts.get("surfaces", 0) + 1
            continue

        if t == "freeform":
            shape = _add_local_paths(slide, x0, y0, bw, bh, el["bbox"],
                                     el.get("paths") or [])
            if shape is not None:
                shape_map[el["id"]] = shape
                counts["freeforms"] = counts.get("freeforms", 0) + 1
            continue

        if t == "raster_crop":
            if original is None:
                original = Image.open(ir["image"]["path"]).convert("RGB")
                assets.mkdir(parents=True, exist_ok=True)
            crop_path = assets / f"{el['id']}.png"
            faithful_crop(original, el, ir["elements"]).save(crop_path)
            pic = slide.shapes.add_picture(str(crop_path), x0, y0, bw, bh)
            shape_map[el["id"]] = pic
            counts["pictures"] += 1
            continue

        if t == "text":
            box = slide.shapes.add_textbox(x0, y0, bw, bh)
            _set_text(box, el, pt_per_px)
            if el.get("rotation"):
                box.rotation = float(el["rotation"])
            shape_map[el["id"]] = box
            counts["shapes"] += 1
            continue

        shape = slide.shapes.add_shape(
            SHAPE_TYPE_MAP.get(t, SHAPE_TYPE_MAP["rect"]), x0, y0, bw, bh)
        if t == "rounded_rect" and el.get("corner") is not None:
            try:
                shape.adjustments[0] = float(el["corner"])
            except Exception:
                pass
        fill = _hex_to_rgb(el.get("fill"))
        if fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
        else:
            shape.fill.background()
        border = _hex_to_rgb(el.get("border_color"))
        if border:
            shape.line.color.rgb = border
            shape.line.width = Pt(el.get("border_width", 1))
            if el.get("dash"):
                from pptx.oxml.ns import qn
                ln = shape.line._get_or_add_ln()
                ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
        else:
            shape.line.fill.background()
        if el.get("text"):
            icon_band = _has_icon_in_top(el, ir["elements"])
            _set_text(shape, el, pt_per_px,
                      avail_frac=0.55 if icon_band else 1.0)
            if icon_band:
                # icon-over-label cards: keep the label out of the icon band
                shape.text_frame.margin_top = int((el["bbox"][3] - el["bbox"][1])
                                                  * 0.42 * scale)
        shape_map[el["id"]] = shape
        counts["shapes"] += 1

    for el in ordered:                               # pass 2: connectors
        if el["type"] not in ("arrow", "line"):
            continue
        if _add_connector(slide, el, shape_map, scale):
            counts["connectors"] += 1

    prs.save(str(out))
    return counts


def _has_icon_in_top(el: dict, elements: list) -> bool:
    x0, y0, x1, y1 = el["bbox"]
    for o in elements:
        if o is el or o.get("type") not in ("icon", "dotcloud") or "bbox" not in o:
            continue
        cx = (o["bbox"][0] + o["bbox"][2]) / 2
        cy = (o["bbox"][1] + o["bbox"][3]) / 2
        if x0 <= cx <= x1 and y0 <= cy <= y0 + (y1 - y0) * 0.5:
            return True
    return False


def _typography_contract(el: dict) -> dict:
    return ((el.get("ext") or {}).get("typography") or {})


def _fit_font_px(el: dict, avail_frac: float = 1.0) -> float:
    """Font size (source px) that fits the bbox in BOTH dimensions.

    Height alone is not enough: a tall narrow box would get a huge size that
    PowerPoint then wraps one syllable per line (the giant-letters bug). The
    width constraint assumes ~0.55 average glyph width/em. An OCR/VLM
    font_size only caps the result, it can't push text out of its box.
    `avail_frac` < 1 when an icon band eats part of the box's height.
    """
    x0, y0, x1, y1 = el["bbox"]
    lines = [ln for ln in (el.get("text") or "").split("\n") if ln] or [" "]
    typo = _typography_contract(el)
    h_factor = float(typo.get("fit_height_factor") or (0.90 if len(lines) == 1 else 0.72))
    h_fit = (y1 - y0) * avail_frac / len(lines) * h_factor
    longest = max(len(ln) for ln in lines)
    width_factor = float(typo.get("fit_width_factor") or 0.53)
    w_fit = (x1 - x0) / max(1, longest) / max(0.20, width_factor)
    size = min(h_fit, w_fit)
    fs = el.get("font_size")
    if fs:
        size = min(size, fs)
    return max(5.0, size)


def _set_text(shape, el: dict, pt_per_px: float,
              avail_frac: float = 1.0) -> None:
    from pptx.enum.text import MSO_AUTO_SIZE

    tf = shape.text_frame
    text = el.get("text", "")
    typo = _typography_contract(el)
    # Manual newlines describe the intended template lines.  Letting
    # PowerPoint auto-wrap those lines can split a final glyph onto its own
    # line ("Engineering" -> "Engineerin" / "g").  Components that truly want
    # automatic paragraph wrapping can opt in through the typography contract.
    if "word_wrap" in typo:
        tf.word_wrap = bool(typo.get("word_wrap"))
    else:
        tf.word_wrap = "\n" in text
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    margin_px = typo.get("margin_px")
    if isinstance(margin_px, (list, tuple)) and len(margin_px) == 4:
        mt, mr, mb, ml = [float(v) * pt_per_px for v in margin_px]
        tf.margin_top = Pt(mt)
        tf.margin_right = Pt(mr)
        tf.margin_bottom = Pt(mb)
        tf.margin_left = Pt(ml)
    else:
        tf.margin_top = tf.margin_bottom = Pt(1)
        tf.margin_left = tf.margin_right = Pt(2)
    size_pt = Pt(max(4, _fit_font_px(el, avail_frac) * pt_per_px))
    color = _hex_to_rgb(el.get("text_color"))
    font_name = el.get("font")
    first = True
    for ln in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = _paragraph_align(el)
        if typo.get("line_spacing"):
            try:
                p.line_spacing = float(typo["line_spacing"])
            except Exception:
                pass
        runs = el.get("runs") if "\n" not in text else None
        if isinstance(runs, list) and runs:
            for spec in runs:
                if not isinstance(spec, dict):
                    continue
                run = p.add_run()
                run.text = str(spec.get("text") or "")
                run.font.size = Pt(max(4, float(spec.get("font_size") or _fit_font_px(el, avail_frac)) * pt_per_px))
                run_color = _hex_to_rgb(spec.get("color") or spec.get("text_color")) or color
                if run_color:
                    run.font.color.rgb = run_color
                if spec.get("bold", el.get("bold")):
                    run.font.bold = True
                if spec.get("italic", el.get("italic")):
                    run.font.italic = True
                run_font = spec.get("font") or font_name
                if run_font:
                    _set_run_font_name(run, str(run_font))
                _set_run_baseline(run, spec.get("baseline"))
        else:
            run = p.add_run()
            run.text = ln
            run.font.size = size_pt
            if color:
                run.font.color.rgb = color
            if el.get("bold"):
                run.font.bold = True
            if el.get("italic"):
                run.font.italic = True
            if font_name:
                _set_run_font_name(run, str(font_name))


def _set_run_baseline(run, baseline) -> None:
    if baseline is None:
        return
    try:
        value = int(float(baseline))
    except (TypeError, ValueError):
        return
    try:
        rpr = run._r.get_or_add_rPr()
        rpr.set("baseline", str(value))
    except Exception:
        pass


def _set_run_font_name(run, font_name: str) -> None:
    run.font.name = font_name
    try:
        from pptx.oxml.ns import qn
        rpr = run._r.get_or_add_rPr()
        rfonts = rpr.get_or_add_latin()
        rfonts.set("typeface", font_name)
        rpr.set(qn("a:dirty"), "0")
    except Exception:
        pass


def _paragraph_align(el: dict):
    align = str(el.get("align") or "").lower()
    if align in {"left", "start"}:
        return PP_ALIGN.LEFT
    if align in {"right", "end"}:
        return PP_ALIGN.RIGHT
    return PP_ALIGN.CENTER


M_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
A14_NS = "http://schemas.microsoft.com/office/drawing/2010/main"


def _set_formula(box, el: dict, pt_per_px: float) -> None:
    """Native equation: OMML wrapped in <a14:m> inside the paragraph.

    Double-click-editable in PowerPoint's equation editor. Falls back to the
    LaTeX source as plain text if the OMML is missing or refuses to parse.
    """
    from pptx.enum.text import MSO_ANCHOR

    tf = box.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # multi-line formulas: title lines as plain runs, math lines as OMML —
    # one paragraph each (the joined single-line version crams and overflows)
    if el.get("omml_lines"):
        try:
            from lxml import etree
            from pptx.oxml.ns import qn
            A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
            n = max(1, len(el["omml_lines"]))
            sz_px = el.get("font_size") or \
                max(5.0, (el["bbox"][3] - el["bbox"][1]) / n * 0.55)
            sz_pt = max(5, sz_px * pt_per_px)
            first = True
            for line in el["omml_lines"]:
                p_obj = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p_obj.alignment = PP_ALIGN.CENTER
                if line["kind"] == "text":
                    run = p_obj.add_run()
                    run.text = line["value"]
                    # card titles read at original weight, not math-line size
                    run.font.size = Pt(sz_pt * 1.25)
                    run.font.bold = True
                    if el.get("text_color"):
                        c = _hex_to_rgb(el["text_color"])
                        if c:
                            run.font.color.rgb = c
                    continue
                wrapped = (f'<m:oMathPara xmlns:m="{M_NS}" xmlns:a="{A_NS}">'
                           f'{line["value"]}</m:oMathPara>')
                omath = etree.fromstring(wrapped.encode("utf-8"))
                sz = str(int(sz_pt * 100))
                for r in omath.iter(f"{{{M_NS}}}r"):
                    rpr = r.makeelement(f"{{{A_NS}}}rPr", {"sz": sz})
                    kids = list(r)
                    pos = 1 if (kids and kids[0].tag == f"{{{M_NS}}}rPr") else 0
                    r.insert(pos, rpr)
                holder = etree.SubElement(p_obj._p, f"{{{A14_NS}}}m",
                                          nsmap={"a14": A14_NS})
                holder.append(omath)
            return
        except Exception:
            pass  # fall through to single-omml / latex paths

    omml = el.get("omml") or _structured_formula_to_omml(el) or _plain_text_to_omml(_formula_plain_text(el))
    if omml:
        try:
            from lxml import etree
            from pptx.oxml.ns import qn
            A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
            wrapped = (f'<m:oMathPara xmlns:m="{M_NS}" xmlns:a="{A_NS}">'
                       f'{omml}</m:oMathPara>'
                       if omml.lstrip().startswith("<m:oMath>") else omml)
            omath = etree.fromstring(wrapped.encode("utf-8"))
            # without an explicit size each run renders at the theme default
            # (~18pt) — gigantic next to a 14px source line. sz = centipoints.
            sz = str(int(max(4, _fit_font_px(el) * pt_per_px) * 100))
            for r in omath.iter(f"{{{M_NS}}}r"):
                rpr = r.makeelement(f"{{{A_NS}}}rPr", {"sz": sz})
                kids = list(r)
                pos = 1 if (kids and kids[0].tag == f"{{{M_NS}}}rPr") else 0
                r.insert(pos, rpr)  # schema order: m:rPr, a:rPr, m:t
            p = tf._txBody.find(qn("a:p"))
            holder = etree.SubElement(p, f"{{{A14_NS}}}m", nsmap={"a14": A14_NS})
            holder.append(omath)
            return
        except Exception:
            pass  # fall through to LaTeX text
    p = tf.paragraphs[0]
    p.alignment = _paragraph_align(el)
    run = p.add_run()
    run.text = _formula_plain_text(el)
    _set_run_font_name(run, str(el.get("font") or "Cambria Math"))
    run.font.size = Pt(max(8, _fit_font_px(el) * pt_per_px))
    color = _hex_to_rgb(el.get("text_color"))
    if color:
        run.font.color.rgb = color
    if el.get("bold"):
        run.font.bold = True
    if el.get("italic"):
        run.font.italic = True


def _formula_plain_text(el: dict) -> str:
    """Readable editable fallback when no OMML is available."""
    latex = str(el.get("latex") or "")
    text = str(el.get("text") or "")
    if "langle" in latex and "beta" in latex and "gamma" in latex:
        return "A = |⟨β, γ⟩| / (∥β∥ ∥γ∥) ≈ 1"
    out = latex or text
    replacements = {
        r"\beta": "β",
        r"\gamma": "γ",
        r"\tau": "τ",
        r"\nabla": "∇",
        r"\approx": "≈",
        r"\langle": "⟨",
        r"\rangle": "⟩",
        r"\|": "∥",
        r"\frac": "",
        "{": "",
        "}": "",
        "$": "",
    }
    for old, new in replacements.items():
        out = out.replace(old, new)
    return " ".join(out.split())


def _plain_text_to_omml(text: str) -> str:
    """Wrap a readable formula string as a minimal editable Office math run."""
    safe = html.escape(str(text or ""), quote=False)
    return (
        f'<m:oMath xmlns:m="{M_NS}">'
        f'<m:r><m:t>{safe}</m:t></m:r>'
        f'</m:oMath>'
    )


def _structured_formula_to_omml(el: dict) -> str | None:
    layout = ((el.get("ext") or {}).get("math_layout") or {})
    if layout.get("kind") == "fraction_expr":
        return _fraction_layout_to_omml(layout)
    tokens = layout.get("tokens") or []
    if not isinstance(tokens, list) or not tokens:
        return None
    parts = []
    for tok in tokens:
        if not isinstance(tok, dict):
            continue
        text = str(tok.get("text") or "")
        if not text:
            continue
        if tok.get("accent") == "hat":
            parts.append(_omml_accent(text, "̂"))
        elif tok.get("accent") == "tilde":
            parts.append(_omml_accent(text, "̃"))
        else:
            parts.append(_omml_run(text))
    if not parts:
        return None
    return f'<m:oMath xmlns:m="{M_NS}">{"".join(parts)}</m:oMath>'


def _fraction_layout_to_omml(layout: dict) -> str | None:
    numerator = _omml_parts(layout.get("numerator") or [])
    denominator = _omml_parts(layout.get("denominator") or [])
    if not numerator or not denominator:
        return None
    prefix = _omml_parts(layout.get("prefix") or [])
    suffix = _omml_parts(layout.get("suffix") or [])
    frac = (
        '<m:f>'
        '<m:fPr><m:type m:val="bar"/></m:fPr>'
        f'<m:num>{"".join(numerator)}</m:num>'
        f'<m:den>{"".join(denominator)}</m:den>'
        '</m:f>'
    )
    return (
        f'<m:oMath xmlns:m="{M_NS}">'
        f'{"".join(prefix)}{frac}{"".join(suffix)}'
        f'</m:oMath>'
    )


def _omml_parts(tokens: list) -> list[str]:
    parts: list[str] = []
    for tok in tokens:
        if not isinstance(tok, dict):
            continue
        text = str(tok.get("text") or "")
        if not text:
            continue
        if tok.get("accent") == "hat":
            parts.append(_omml_accent(text, "̂"))
        elif tok.get("accent") == "tilde":
            parts.append(_omml_accent(text, "̃"))
        else:
            parts.append(_omml_run(text))
    return parts


def _omml_run(text: str) -> str:
    safe = html.escape(str(text or ""), quote=False)
    return f"<m:r><m:t>{safe}</m:t></m:r>"


def _omml_hat(text: str) -> str:
    return _omml_accent(text, "̂")


def _omml_accent(text: str, accent: str) -> str:
    safe = html.escape(str(text or ""), quote=False)
    safe_accent = html.escape(str(accent or ""), quote=False)
    return (
        '<m:acc>'
        f'<m:accPr><m:chr m:val="{safe_accent}"/></m:accPr>'
        f'<m:e><m:r><m:t>{safe}</m:t></m:r></m:e>'
        '</m:acc>'
    )


def _add_chart(slide, el: dict, x0: int, y0: int, bw: int, bh: int):
    """Native PPT chart from the VLM-extracted spec (values are approximate
    — ext.approx is set; the point is the data is now EDITABLE)."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    paths = el.get("paths") or (el.get("chart") or {}).get("paths") \
        or (el.get("ext", {}).get("chart") or {}).get("paths")
    if paths:
        rendered = _add_local_paths(slide, x0, y0, bw, bh, el["bbox"], paths)
        if rendered is not None:
            return rendered

    spec = _normalize_chart_spec(el.get("chart") or el.get("ext", {}).get("chart") or {})
    try:
        data = CategoryChartData()
        data.categories = [str(c) for c in spec["categories"]]
        for s in spec["series"]:
            data.add_series(s.get("name") or "series", s["values"])
        ctype = (XL_CHART_TYPE.COLUMN_CLUSTERED if spec["kind"] == "bar"
                 else XL_CHART_TYPE.LINE)
        frame = slide.shapes.add_chart(ctype, x0, y0, bw, bh, data)
        chart = frame.chart
        chart.has_title = False
        chart.has_legend = len(spec["series"]) > 1
        if chart.has_legend:
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(7)
        # default chart text (~18pt) dwarfs a 2-inch chart
        for axis in (chart.category_axis, chart.value_axis):
            try:
                axis.tick_labels.font.size = Pt(7)
            except Exception:
                pass
        for i, s in enumerate(spec["series"]):
            color = _hex_to_rgb(s.get("color"))
            if color is None:
                continue
            srs = chart.series[i]
            srs.format.fill.solid()
            srs.format.fill.fore_color.rgb = color
            if spec["kind"] == "line":
                srs.format.line.color.rgb = color
        return frame
    except Exception:
        return None


def _normalize_chart_spec(spec: dict) -> dict:
    """Accept v2/v3 chart variants and return the schema python-pptx needs."""
    if not isinstance(spec, dict):
        return {"kind": "none", "categories": [], "series": []}

    kind = str(spec.get("kind") or spec.get("type") or "none").lower()
    if kind not in ("bar", "line", "scatter", "pie"):
        kind = "line" if spec.get("series") else "none"

    series_out = []
    categories = list(spec.get("categories") or [])
    for s in spec.get("series") or []:
        if not isinstance(s, dict):
            continue
        values = list(s.get("values") or [])
        points = s.get("points") or []
        if not values and points:
            values = [float(p.get("y", 0.0)) for p in points if isinstance(p, dict)]
            if not categories:
                categories = [str(i + 1) for i in range(len(values))]
        if values:
            series_out.append({
                "name": str(s.get("name") or "series"),
                "color": s.get("color") or "#4472c4",
                "values": values,
            })

    if series_out and (not categories or len(categories) != len(series_out[0]["values"])):
        categories = [str(i + 1) for i in range(len(series_out[0]["values"]))]

    return {"kind": kind, "categories": categories, "series": series_out}


ICON_SHAPE_MAP = {
    "database": "CAN",
    "gear": "GEAR_6",
    "warning": "ISOSCELES_TRIANGLE",
    "arrow": "RIGHT_ARROW",
}


def _add_icon(slide, el: dict, x0: int, y0: int, bw: int, bh: int,
              pt_per_px: float):
    """Pictogram → MSO autoshape when one matches, else a glyph textbox."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    icon = el.get("icon") or {}
    kind = icon.get("kind", "other")
    color = _hex_to_rgb(icon.get("color"))
    if icon.get("paths"):
        last = _add_local_paths(slide, x0, y0, bw, bh, el["bbox"],
                                icon["paths"])
        if last is not None:
            return last

    if kind == "chart":          # three mini bars beat an emoji glyph
        heights = (0.55, 0.95, 0.75)
        bar_w = bw // 4
        for i, hf in enumerate(heights):
            bar = slide.shapes.add_shape(
                SHAPE_TYPE_MAP["rect"], x0 + i * (bar_w + bar_w // 4),
                y0 + int(bh * (1 - hf)), bar_w, int(bh * hf))
            bar.fill.solid()
            bar.fill.fore_color.rgb = color or _hex_to_rgb("#4472C4")
            bar.line.fill.background()
        return bar

    if kind == "hourglass":
        return _add_hourglass_icon(slide, x0, y0, bw, bh,
                                   color or _hex_to_rgb("#9a5b13"))

    if kind == "document":
        return _add_document_icon(slide, x0, y0, bw, bh,
                                  color or _hex_to_rgb("#245591"))

    if kind == "warning" and icon.get("variant") == "outline":
        return _add_warning_outline_icon(slide, x0, y0, bw, bh,
                                         color or _hex_to_rgb("#ef4b36"),
                                         pt_per_px)

    if kind == "eye_slash":
        return _add_eye_slash_icon(slide, x0, y0, bw, bh,
                                   color or _hex_to_rgb("#4b718f"))

    if kind == "bell":           # distribution curve (the CI-estimator icon)
        import math as _m
        from pptx.util import Emu
        pts = [(t, _m.exp(-((t - 0.5) ** 2) / 0.045)) for t in
               [i / 30 for i in range(31)]]
        px0, py0, pw, ph = el["bbox"][0], el["bbox"][1], \
            el["bbox"][2] - el["bbox"][0], el["bbox"][3] - el["bbox"][1]
        scale_l = bw / max(1e-6, pw)   # EMU per source px
        fb = slide.shapes.build_freeform(pts[0][0] * pw, (1 - pts[0][1]) * ph,
                                         scale=scale_l)
        fb.add_line_segments([(t * pw, (1 - v) * ph) for t, v in pts[1:]],
                             close=False)
        curve = fb.convert_to_shape(origin_x=Emu(x0), origin_y=Emu(y0))
        curve.fill.background()
        curve.line.color.rgb = color or _hex_to_rgb("#555555")
        curve.line.width = Pt(1.4)
        return curve

    if kind in ("line", "scatter"):
        return _add_line_or_scatter_icon(slide, el, x0, y0, bw, bh,
                                         color or _hex_to_rgb("#4f7db8"),
                                         with_points=(kind == "scatter"))

    if kind == "shield":
        return _add_shield_icon(slide, x0, y0, bw, bh,
                                color or _hex_to_rgb("#3aa27e"))

    if kind == "check":
        return _add_check_icon(slide, x0, y0, bw, bh,
                               color or _hex_to_rgb("#2b9c6a"))

    if kind == "cross":
        return _add_cross_icon(slide, x0, y0, bw, bh,
                               color or _hex_to_rgb("#c94235"))

    mso_name = ICON_SHAPE_MAP.get(kind)
    if mso_name:
        shape = slide.shapes.add_shape(getattr(MSO_SHAPE, mso_name),
                                       x0, y0, bw, bh)
        if color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        if kind == "warning":
            tf = shape.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            run = tf.paragraphs[0].add_run()
            run.text = "!"
            run.font.bold = True
            run.font.size = Pt(max(6, (el["bbox"][3] - el["bbox"][1]) * 0.5 * pt_per_px))
        return shape
    if kind == "other":
        return _add_generic_icon(slide, x0, y0, bw, bh,
                                 color or _hex_to_rgb("#4f5d73"))
    box = slide.shapes.add_textbox(x0, y0, bw, bh)
    tf = box.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = icon.get("glyph") or "◆"
    run.font.size = Pt(max(6, (el["bbox"][3] - el["bbox"][1]) * 0.7 * pt_per_px))
    if color:
        run.font.color.rgb = color
    return box


def _add_generic_icon(slide, x0: int, y0: int, bw: int, bh: int, color):
    """Fallback icon as native geometry, not a text glyph."""
    from pptx.enum.shapes import MSO_SHAPE

    size = max(2, min(bw, bh) // 5)
    positions = [
        (0.25, 0.30),
        (0.50, 0.48),
        (0.75, 0.28),
        (0.35, 0.72),
        (0.70, 0.68),
    ]
    last = None
    for fx, fy in positions:
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x0 + int(fx * bw) - size // 2,
            y0 + int(fy * bh) - size // 2,
            size,
            size,
        )
        if color:
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        last = dot
    return last


def _add_hourglass_icon(slide, x0: int, y0: int, bw: int, bh: int, color):
    """Outlined hourglass icon as native editable line geometry."""
    from pptx.util import Emu

    cx = x0 + bw * 0.5
    top_y = y0 + bh * 0.12
    bot_y = y0 + bh * 0.88
    neck_y = y0 + bh * 0.50
    left = x0 + bw * 0.22
    right = x0 + bw * 0.78
    neck_l = x0 + bw * 0.43
    neck_r = x0 + bw * 0.57

    paths = [
        [(left, top_y), (right, top_y)],
        [(left, bot_y), (right, bot_y)],
        [(left, top_y), (neck_l, neck_y), (left, bot_y)],
        [(right, top_y), (neck_r, neck_y), (right, bot_y)],
    ]
    last = None
    for pts in paths:
        fb = slide.shapes.build_freeform(pts[0][0] - x0, pts[0][1] - y0, scale=1.0)
        fb.add_line_segments([(px - x0, py - y0) for px, py in pts[1:]], close=False)
        shp = fb.convert_to_shape(origin_x=Emu(x0), origin_y=Emu(y0))
        shp.fill.background()
        shp.line.color.rgb = color
        shp.line.width = Pt(2.1)
        last = shp

    for yy, ww in ((y0 + bh * 0.36, bw * 0.20), (y0 + bh * 0.64, bw * 0.25)):
        sand = slide.shapes.add_shape(
            SHAPE_TYPE_MAP["oval"],
            int(cx - ww / 2),
            int(yy - bh * 0.025),
            int(ww),
            int(bh * 0.05),
        )
        sand.fill.solid()
        sand.fill.fore_color.rgb = color
        sand.line.fill.background()
        last = sand
    return last


def _add_document_icon(slide, x0: int, y0: int, bw: int, bh: int, color):
    """Outlined folded-corner report page with native line details."""
    from pptx.util import Emu

    left = x0 + bw * 0.18
    top = y0 + bh * 0.10
    right = x0 + bw * 0.82
    bottom = y0 + bh * 0.90
    fold = min(bw, bh) * 0.22
    pts = [
        (left, top),
        (right - fold, top),
        (right, top + fold),
        (right, bottom),
        (left, bottom),
        (left, top),
    ]
    fb = slide.shapes.build_freeform(pts[0][0] - x0, pts[0][1] - y0, scale=1.0)
    fb.add_line_segments([(px - x0, py - y0) for px, py in pts[1:]], close=False)
    page = fb.convert_to_shape(origin_x=Emu(x0), origin_y=Emu(y0))
    page.fill.solid()
    page.fill.fore_color.rgb = _hex_to_rgb("#fbfdff")
    page.line.color.rgb = color
    page.line.width = Pt(1.6)

    fold_line = _ppt_line(slide, right - fold, top, right - fold, top + fold, color, 1.2)
    _ppt_line(slide, right - fold, top + fold, right, top + fold, color, 1.2)

    for i, frac in enumerate((0.34, 0.47, 0.60)):
        _ppt_line(
            slide,
            left + bw * 0.12,
            top + bh * frac,
            right - bw * 0.16,
            top + bh * frac,
            color,
            1.15,
        )
    bars = [(0.24, 0.20), (0.40, 0.33), (0.56, 0.45)]
    last = fold_line
    for i, (fx, fh) in enumerate(bars):
        bar = slide.shapes.add_shape(
            SHAPE_TYPE_MAP["rect"],
            int(left + bw * (0.13 + i * 0.13)),
            int(bottom - bh * (0.12 + fh)),
            int(bw * 0.08),
            int(bh * fh),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        last = bar
    return last


def _add_warning_outline_icon(slide, x0: int, y0: int, bw: int, bh: int,
                              color, pt_per_px: float):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, x0, y0, bw, bh)
    tri.fill.background()
    tri.line.color.rgb = color
    tri.line.width = Pt(2.0)

    box = slide.shapes.add_textbox(
        int(x0 + bw * 0.30),
        int(y0 + bh * 0.18),
        int(bw * 0.40),
        int(bh * 0.58),
    )
    tf = box.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "!"
    run.font.bold = True
    run.font.size = Pt(max(6, (min(bw, bh) / 12700.0) * 0.42))
    run.font.color.rgb = color
    return box


def _add_eye_slash_icon(slide, x0: int, y0: int, bw: int, bh: int, color):
    eye = slide.shapes.add_shape(
        SHAPE_TYPE_MAP["oval"],
        int(x0 + bw * 0.10),
        int(y0 + bh * 0.22),
        int(bw * 0.80),
        int(bh * 0.50),
    )
    eye.fill.background()
    eye.line.color.rgb = color
    eye.line.width = Pt(1.8)
    pupil = slide.shapes.add_shape(
        SHAPE_TYPE_MAP["oval"],
        int(x0 + bw * 0.43),
        int(y0 + bh * 0.39),
        int(bw * 0.14),
        int(bh * 0.14),
    )
    pupil.fill.solid()
    pupil.fill.fore_color.rgb = color
    pupil.line.fill.background()
    slash = _ppt_line(
        slide,
        x0 + bw * 0.15,
        y0 + bh * 0.82,
        x0 + bw * 0.88,
        y0 + bh * 0.12,
        color,
        2.2,
    )
    return slash


def _ppt_line(slide, x0, y0, x1, y1, color, width_pt: float):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, int(x0), int(y0), int(x1), int(y1))
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    return line


def _add_local_paths(slide, x0: int, y0: int, bw: int, bh: int,
                     bbox: list, paths: list[dict]):
    """Render crop-local contours as editable PowerPoint freeforms."""
    from pptx.util import Emu

    src_w = max(1.0, float(bbox[2]) - float(bbox[0]))
    src_h = max(1.0, float(bbox[3]) - float(bbox[1]))
    scale_l = min(bw / src_w, bh / src_h)
    last = None
    for path in sorted(paths, key=lambda p: p.get("area", 0), reverse=True):
        pts = path.get("points") or []
        closed = bool(path.get("closed", True))
        if len(pts) < (3 if closed else 2):
            continue
        fb = slide.shapes.build_freeform(pts[0][0], pts[0][1],
                                         scale=scale_l)
        fb.add_line_segments([(p[0], p[1]) for p in pts[1:]], close=closed)
        shape = fb.convert_to_shape(origin_x=Emu(x0), origin_y=Emu(y0))
        fill = _hex_to_rgb(path.get("fill"))
        if fill and closed:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
            alpha = path.get("alpha")
            if alpha is not None:
                _set_fill_alpha(shape, int(alpha))
        else:
            shape.fill.background()
        line = _hex_to_rgb(path.get("line"))
        if line:
            shape.line.color.rgb = line
            shape.line.width = Pt(float(path.get("line_width", 0.4)))
            alpha = path.get("alpha")
            if alpha is not None:
                _set_line_alpha(shape, int(alpha))
        else:
            shape.line.fill.background()
        last = shape
    return last


def _add_line_or_scatter_icon(slide, el: dict, x0: int, y0: int, bw: int,
                              bh: int, color, with_points: bool):
    import math as _m
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    pts = []
    for i in range(18):
        t = i / 17
        y = 0.62 - 0.34 * t + 0.10 * _m.sin(t * _m.pi * 2.0)
        pts.append((t, max(0.1, min(0.9, y))))
    fb = slide.shapes.build_freeform(pts[0][0] * bw, pts[0][1] * bh, scale=1.0)
    fb.add_line_segments([(t * bw, y * bh) for t, y in pts[1:]], close=False)
    curve = fb.convert_to_shape(origin_x=Emu(x0), origin_y=Emu(y0))
    curve.fill.background()
    if color:
        curve.line.color.rgb = color
    curve.line.width = Pt(1.3)
    if not with_points:
        return curve

    for i, (t, y) in enumerate(pts[::3]):
        r = max(2, min(bw, bh) // 18)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x0 + int(t * bw) - r,
            y0 + int(y * bh) - r,
            2 * r,
            2 * r,
        )
        if color:
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
        dot.line.fill.background()
    return curve


def _add_shield_icon(slide, x0: int, y0: int, bw: int, bh: int, color):
    from pptx.util import Emu

    pts = [
        (0.50 * bw, 0.05 * bh),
        (0.88 * bw, 0.20 * bh),
        (0.80 * bw, 0.66 * bh),
        (0.50 * bw, 0.95 * bh),
        (0.20 * bw, 0.66 * bh),
        (0.12 * bw, 0.20 * bh),
    ]
    fb = slide.shapes.build_freeform(pts[0][0], pts[0][1], scale=1.0)
    fb.add_line_segments(pts[1:], close=True)
    shield = fb.convert_to_shape(origin_x=Emu(x0), origin_y=Emu(y0))
    if color:
        shield.fill.solid()
        shield.fill.fore_color.rgb = color
    shield.line.fill.background()

    chk = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        x0 + int(0.32 * bw), y0 + int(0.52 * bh),
        x0 + int(0.46 * bw), y0 + int(0.67 * bh),
    )
    chk.line.color.rgb = _hex_to_rgb("#ffffff")
    chk.line.width = Pt(2.0)
    chk2 = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        x0 + int(0.46 * bw), y0 + int(0.67 * bh),
        x0 + int(0.70 * bw), y0 + int(0.36 * bh),
    )
    chk2.line.color.rgb = _hex_to_rgb("#ffffff")
    chk2.line.width = Pt(2.0)
    return shield


def _add_check_icon(slide, x0: int, y0: int, bw: int, bh: int, color):
    a = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        x0 + int(0.20 * bw), y0 + int(0.55 * bh),
        x0 + int(0.42 * bw), y0 + int(0.75 * bh),
    )
    b = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        x0 + int(0.42 * bw), y0 + int(0.75 * bh),
        x0 + int(0.82 * bw), y0 + int(0.25 * bh),
    )
    for s in (a, b):
        if color:
            s.line.color.rgb = color
        s.line.width = Pt(2.2)
    return b


def _add_cross_icon(slide, x0: int, y0: int, bw: int, bh: int, color):
    a = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        x0 + int(0.20 * bw), y0 + int(0.20 * bh),
        x0 + int(0.80 * bw), y0 + int(0.80 * bh),
    )
    b = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        x0 + int(0.80 * bw), y0 + int(0.20 * bh),
        x0 + int(0.20 * bw), y0 + int(0.80 * bh),
    )
    for s in (a, b):
        if color:
            s.line.color.rgb = color
        s.line.width = Pt(2.2)
    return b


def _add_dotcloud(slide, el: dict, scale: float) -> int:
    """Organic art: smooth gradient surface + flow lines + scatter dots."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    from work.diagram2ppt.v2.vectorize import smooth_closed

    bx0, by0 = el["bbox"][0], el["bbox"][1]
    path_shapes = 0
    if el.get("paths"):
        rendered = _add_local_paths(
            slide,
            int(bx0 * scale),
            int(by0 * scale),
            int((el["bbox"][2] - el["bbox"][0]) * scale),
            int((el["bbox"][3] - el["bbox"][1]) * scale),
            el["bbox"],
            el["paths"],
        )
        if rendered is not None:
            path_shapes = len(el.get("paths") or [])

    # wave bands stand alone — they replace the silhouette entirely, so they
    # must NOT live inside the silhouette branch (the v62 invisible-surface bug)
    wb = el.get("wave_bands")
    if wb and len(wb.get("curves", [])) >= 2:
        from pptx.util import Emu as _Emu
        curves, fills = wb["curves"], wb.get("fills") or []
        for bi in range(len(curves) - 1):
            a, b = curves[bi], curves[bi + 1]
            path = a + b[::-1]
            fb2 = slide.shapes.build_freeform(path[0][0], path[0][1],
                                              scale=scale)
            fb2.add_line_segments([(p[0], p[1]) for p in path[1:]], close=True)
            band = fb2.convert_to_shape(origin_x=_Emu(int(bx0 * scale)),
                                        origin_y=_Emu(int(by0 * scale)))
            fill_value = fills[bi] if bi < len(fills) else (fills[-1] if fills else "#dbeaf2")
            col = _hex_to_rgb(fill_value)
            if col:
                band.fill.solid()
                band.fill.fore_color.rgb = col
                # the surface fades downward into a veil — a solid slab at the
                # bottom was the m3 native-res finding
                fade = [92, 76, 55, 34]
                _set_fill_alpha(band, fade[min(bi, len(fade) - 1)])
            band.line.fill.background()
            _set_soft_edge(band, 2.5)
        from pptx.dml.color import RGBColor
        for ci, c in enumerate(curves[1:-1], 1):       # ridge lines
            lb2 = slide.shapes.build_freeform(c[0][0], c[0][1], scale=scale)
            lb2.add_line_segments([(p[0], p[1]) for p in c[1:]], close=False)
            ridge = lb2.convert_to_shape(origin_x=_Emu(int(bx0 * scale)),
                                         origin_y=_Emu(int(by0 * scale)))
            ridge.fill.background()
            rc = _hex_to_rgb(fills[min(ci, len(fills) - 1)])
            if rc:
                hx = str(rc)
                ridge.line.color.rgb = RGBColor(
                    *(max(0, int(hx[i:i + 2], 16) - 45) for i in (0, 2, 4)))
            ridge.line.width = Pt(0.6)
            _set_line_alpha(ridge, 60)

    sil = el.get("silhouette")
    if sil and len(sil.get("points", [])) >= 3:
        pts = smooth_closed(sil["points"])     # dense curve, no spikes
        fb = slide.shapes.build_freeform(pts[0][0], pts[0][1], scale=scale)
        fb.add_line_segments([(p[0], p[1]) for p in pts[1:]], close=True)
        shape = fb.convert_to_shape(origin_x=Emu(int(bx0 * scale)),
                                    origin_y=Emu(int(by0 * scale)))
        style = el.get("style") or {}
        light = _hex_to_rgb(style.get("light") or sil.get("fill"))
        dark = _hex_to_rgb(style.get("dark"))
        shape.line.fill.background()
        if light and dark:
            shape.fill.gradient()              # soft shading sells the 3D
            stops = shape.fill.gradient_stops
            stops[0].color.rgb = light
            stops[0].position = 0.0
            stops[1].color.rgb = dark
            stops[1].position = 1.0
            try:
                shape.fill.gradient_angle = 115.0
            except Exception:
                pass
        elif light:
            shape.fill.solid()
            shape.fill.fore_color.rgb = light
        # relief bands back-to-front, progressively more translucent with
        # native soft edges — airiness, not terraces
        n_layers = max(1, len(el.get("surface_layers", [])))
        for li, layer in enumerate(el.get("surface_layers", [])):
            lcol = _hex_to_rgb(layer["fill"])
            for poly in layer["polys"]:
                pts2 = smooth_closed(poly, samples=110)
                pb = slide.shapes.build_freeform(pts2[0][0], pts2[0][1],
                                                 scale=scale)
                pb.add_line_segments([(p[0], p[1]) for p in pts2[1:]],
                                     close=True)
                ply = pb.convert_to_shape(origin_x=Emu(int(bx0 * scale)),
                                          origin_y=Emu(int(by0 * scale)))
                if lcol:
                    ply.fill.solid()
                    ply.fill.fore_color.rgb = lcol
                    _set_fill_alpha(ply, 88 - int(26 * li / n_layers))
                if li == n_layers - 1:      # crest highlight along the ridge
                    from pptx.dml.color import RGBColor
                    ply.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    ply.line.width = Pt(0.75)
                    _set_line_alpha(ply, 70)
                else:
                    ply.line.fill.background()
                _set_soft_edge(ply, 4.0)

        pass
    for heat in el.get("heat_regions", []) or []:
        try:
            cx = float(heat.get("cx", 0.0))
            cy = float(heat.get("cy", 0.0))
            rx = float(heat.get("rx", 0.0))
            ry = float(heat.get("ry", 0.0))
        except Exception:
            continue
        if rx <= 1 or ry <= 1:
            continue
        ox = int((bx0 + cx - rx) * scale)
        oy = int((by0 + cy - ry) * scale)
        ow = max(1, int(2 * rx * scale))
        oh = max(1, int(2 * ry * scale))
        blob = slide.shapes.add_shape(MSO_SHAPE.OVAL, ox, oy, ow, oh)
        color = _hex_to_rgb(heat.get("color") or "#d8e8f3")
        if color:
            blob.fill.solid()
            blob.fill.fore_color.rgb = color
            _set_fill_alpha(blob, int(heat.get("opacity", 40)))
        blob.line.fill.background()
        _set_soft_edge(blob, float(heat.get("soft_edge", 8.0)))
    # fabric-of-space contour lines — independent of any silhouette (twice
    # now a feature died silently inside that branch)
    sl_dark = _hex_to_rgb((el.get("style") or {}).get("dark") or "#b6c2d4")
    for line in el.get("streamlines", []):
        if len(line) < 3:
            continue
        lb = slide.shapes.build_freeform(line[0][0], line[0][1], scale=scale)
        lb.add_line_segments([(p[0], p[1]) for p in line[1:]], close=False)
        ln = lb.convert_to_shape(origin_x=Emu(int(bx0 * scale)),
                                 origin_y=Emu(int(by0 * scale)))
        ln.fill.background()
        if sl_dark:
            ln.line.color.rgb = sl_dark
        ln.line.width = Pt(0.5)
        _set_line_alpha(ln, 50)

    n = path_shapes
    dots = el.get("dots", [])
    clip_poly = None
    if wb and len(wb.get("curves", [])) >= 2:
        # third silently-dead clip: the old guard referenced the silhouette
        # field that wave_bands replaced — re-anchor to the wave envelope
        clip_poly = wb["curves"][0] + wb["curves"][-1][::-1]
    elif sil and len(sil.get("points", [])) >= 3:
        clip_poly = sil["points"]
    if clip_poly:
        from matplotlib.path import Path as MplPath
        poly = MplPath(clip_poly)
        dots = [d for d in dots if poly.contains_point((d["cx"], d["cy"]),
                                                       radius=10)]
    bh_el = max(1.0, el["bbox"][3] - el["bbox"][1])
    big_field = len(dots) > 40            # depth fade only on large surfaces
    for d in dots:
        r = max(2.2, d["r"])   # sub-2px dots vanish at slide scale
        ox = int((bx0 + d["cx"] - r) * scale)
        oy = int((by0 + d["cy"] - r) * scale)
        size = max(1, int(2 * r * scale))
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, ox, oy, size, size)
        color = _hex_to_rgb(d.get("color"))
        if color:
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            if big_field:                  # far (top) dots fade — depth cue
                depth = 1.0 - d["cy"] / bh_el        # 1=front(bottom) 0=far
                _set_fill_alpha(dot, int(55 + 45 * depth))
        dot.line.fill.background()
        n += 1
    return n


def _set_soft_edge(shape, radius_pt: float) -> None:
    """Native PowerPoint feathered edge (a:softEdge) — the airiness the
    hard-edged bands lacked. radius in points."""
    try:
        from pptx.oxml.ns import qn
        spPr = shape._element.spPr
        eff = spPr.find(qn("a:effectLst"))
        if eff is None:
            eff = spPr.makeelement(qn("a:effectLst"), {})
            spPr.append(eff)
        eff.append(eff.makeelement(qn("a:softEdge"),
                                   {"rad": str(int(radius_pt * 12700))}))
    except Exception:
        pass


def _set_line_alpha(shape, opacity_pct: int) -> None:
    try:
        from pptx.oxml.ns import qn
        ln = shape.line._get_or_add_ln()
        srgb = ln.find(qn("a:solidFill") + "/" + qn("a:srgbClr"))
        if srgb is None:
            srgb = ln.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        srgb.append(srgb.makeelement(qn("a:alpha"),
                                     {"val": str(int(opacity_pct * 1000))}))
    except Exception:
        pass


def _set_fill_alpha(shape, opacity_pct: int) -> None:
    """Translucency lets lower relief layers show through (depth cue)."""
    try:
        from pptx.oxml.ns import qn
        srgb = shape.fill._xPr.find(qn("a:solidFill") + "/" + qn("a:srgbClr"))
        if srgb is None:
            srgb = shape.fill._xPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        alpha = srgb.makeelement(qn("a:alpha"),
                                 {"val": str(int(opacity_pct * 1000))})
        srgb.append(alpha)
    except Exception:
        pass


def _add_connector(slide, el: dict, shape_map: dict, scale: float) -> bool:
    src = shape_map.get(el.get("from_id") or "")
    dst = shape_map.get(el.get("to_id") or "")
    if src is not None and dst is not None:
        scx, scy = src.left + src.width // 2, src.top + src.height // 2
        dcx, dcy = dst.left + dst.width // 2, dst.top + dst.height // 2
        start = _edge(src, (dcx, dcy))
        end = _edge(dst, (scx, scy))
    elif el.get("points"):
        p = el["points"]
        start = (int(p[0] * scale), int(p[1] * scale))
        end = (int(p[2] * scale), int(p[3] * scale))
        if el.get("thickness", 0) >= 6:   # fat stroke = block arrow, not a hairline
            return _add_block_arrow(slide, el, start, end, scale)
    elif el.get("start") and el.get("end"):
        p0, p1 = el["start"], el["end"]
        start = (int(float(p0[0]) * scale), int(float(p0[1]) * scale))
        end = (int(float(p1[0]) * scale), int(float(p1[1]) * scale))
        if el.get("thickness", 0) >= 6:
            return _add_block_arrow(slide, el, start, end, scale)
    else:
        return False
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, start[0], start[1], end[0], end[1])
    color = _hex_to_rgb(el.get("color"))
    if color:
        conn.line.color.rgb = color
    conn.line.width = Pt(float(el.get("line_width") or el.get("thickness") or 1.5))
    if el.get("dash"):
        from pptx.oxml.ns import qn
        ln = conn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    if el["type"] == "arrow":
        _add_arrowhead(conn)
    return True


def _add_block_arrow(slide, el: dict, start: tuple, end: tuple,
                     scale: float) -> bool:
    """Fat directional stroke → rotated RIGHT_ARROW autoshape."""
    import math
    from pptx.enum.shapes import MSO_SHAPE

    dx, dy = end[0] - start[0], end[1] - start[1]
    L = int(math.hypot(dx, dy))
    if L < 1:
        return False
    H = int(min(44, max(6, el.get("thickness", 8))) * 1.5 * scale)
    cx, cy = (start[0] + end[0]) // 2, (start[1] + end[1]) // 2
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   cx - L // 2, cy - H // 2, L, H)
    shape.rotation = math.degrees(math.atan2(dy, dx))
    color = _hex_to_rgb(el.get("color"))
    if color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return True


def _edge(shape, toward: tuple) -> tuple:
    cx, cy = shape.left + shape.width // 2, shape.top + shape.height // 2
    dx, dy = toward[0] - cx, toward[1] - cy
    if dx == 0 and dy == 0:
        return (cx, cy)
    hw, hh = shape.width / 2, shape.height / 2
    sx = hw / abs(dx) if dx else float("inf")
    sy = hh / abs(dy) if dy else float("inf")
    s = min(sx, sy)
    return (int(cx + dx * s), int(cy + dy * s))
